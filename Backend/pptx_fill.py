from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import re
import tempfile
from drive_utils import get_drive_service, download_file, update_file
from combined import build_full_data

# Extra keys so older templates with shortened tag names still get GA4 property totals.
_TAG_ALIASES = {
    "ga4_users": "ga4_total_users",
    "ga4_sessions": "ga4_total_sessions",
    "ga4_views": "ga4_total_views",
    "ga4_duration": "ga4_total_duration",
    "ga4_bounce": "ga4_total_bounce",
    "ga4_buttons": "ga4_total_buttons",
    "total_users": "ga4_total_users",
    "total_sessions": "ga4_total_sessions",
    "total_views": "ga4_total_views",
    "ga4_users_prev": "ga4_total_users_prev",
    "ga4_sessions_prev": "ga4_total_sessions_prev",
    "ga4_views_prev": "ga4_total_views_prev",
    "ga4_users_pct": "ga4_total_users_pct",
    "ga4_sessions_pct": "ga4_total_sessions_pct",
    "ga4_views_pct": "ga4_total_views_pct",
}

# Tags that must be replaced for the main GA4 summary table to match combined.py / GA4.
_CRITICAL_GA4_TAGS = (
    "ga4_total_views",
    "ga4_total_sessions",
    "ga4_total_users",
    "ga4_total_duration",
    "ga4_total_bounce",
    "ga4_total_views_prev",
    "ga4_total_sessions_prev",
    "ga4_total_users_prev",
)

_TAG_IN_TEXT_RE = re.compile(r"\{\{([^}]+)\}\}|PLACEHOLDER_([A-Za-z0-9_]+)")


def _expand_data_aliases(data: dict) -> dict:
    merged = dict(data)
    for alias, canonical in _TAG_ALIASES.items():
        if canonical in data and alias not in merged:
            merged[alias] = data[canonical]
    return merged


def _sorted_replacement_items(data: dict):
    """Longest keys first so {{ga4_total_views_pct}} is not broken by {{ga4_total_views}}."""
    return sorted(data.items(), key=lambda item: len(item[0]), reverse=True)


def apply_tag_replacements(text: str, data: dict, found_list: list) -> str:
    if not text or ("{{" not in text and "PLACEHOLDER_" not in text):
        return text
    out = text
    for key, value in _sorted_replacement_items(data):
        tag1 = f"{{{{{key}}}}}"
        tag2 = f"PLACEHOLDER_{key}"
        if tag1 in out:
            out = out.replace(tag1, str(value))
            found_list.append(key)
        if tag2 in out:
            out = out.replace(tag2, str(value))
            found_list.append(key)
    return out


def replace_in_paragraphs(paragraphs, data, slide_num, missing_list, found_list):
    for para in paragraphs:
        full_para_text = "".join(run.text for run in para.runs)
        original_text = full_para_text

        if "{{" not in full_para_text and "PLACEHOLDER_" not in full_para_text:
            continue

        full_para_text = apply_tag_replacements(full_para_text, data, found_list)

        if full_para_text != original_text:
            if para.runs:
                para.runs[0].text = full_para_text
                for run in para.runs[1:]:
                    run.text = ""

        if "{{" in full_para_text or "PLACEHOLDER_" in full_para_text:
            missing_list.append(f"Slide {slide_num}: {full_para_text[:80]}")


def replace_in_text_frame(text_frame, data, slide_num, missing_list, found_list):
    if text_frame is None:
        return
    replace_in_paragraphs(text_frame.paragraphs, data, slide_num, missing_list, found_list)


def replace_in_table(table, data, slide_num, missing_list, found_list):
    for row in table.rows:
        for cell in row.cells:
            raw = cell.text or ""
            if "{{" in raw or "PLACEHOLDER_" in raw:
                updated = apply_tag_replacements(raw, data, found_list)
                if updated != raw:
                    cell.text = updated
                    raw = updated
            replace_in_text_frame(cell.text_frame, data, slide_num, missing_list, found_list)


def iter_shapes(shapes):
    """Yield all shapes including those nested inside GROUP shapes."""
    for shape in shapes:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def process_shape(shape, data, slide_num, missing_list, found_list):
    if getattr(shape, "has_table", False) and shape.has_table:
        replace_in_table(shape.table, data, slide_num, missing_list, found_list)
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        replace_in_text_frame(shape.text_frame, data, slide_num, missing_list, found_list)


def collect_template_tags(prs: Presentation) -> set[str]:
    tags: set[str] = set()

    def scan_text(text: str):
        if not text:
            return
        for m in _TAG_IN_TEXT_RE.finditer(text):
            tags.add((m.group(1) or m.group(2) or "").strip())

    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        scan_text(cell.text)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    scan_text("".join(run.text for run in para.runs))
    return {t for t in tags if t}


def fill_presentation(data, template_path, output_path):
    """Open template, replace tags in all shapes (including grouped tables), save."""
    data = _expand_data_aliases(data)
    prs = Presentation(template_path)
    placeholders_found: list = []
    placeholders_missing: list = []

    template_tags = collect_template_tags(prs)
    unknown_tags = sorted(
        t for t in template_tags
        if t not in data and t not in _TAG_ALIASES
    )
    if unknown_tags:
        print("WARNING — template contains tags not in build_full_data (cells may keep old numbers):")
        for t in unknown_tags[:25]:
            print(f"  - {{{{{t}}}}}")
        if len(unknown_tags) > 25:
            print(f"  ... and {len(unknown_tags) - 25} more")

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in iter_shapes(slide.shapes):
            process_shape(shape, data, slide_num, placeholders_missing, placeholders_found)

    prs.save(output_path)

    found_set = set(placeholders_found)
    missing_critical = [k for k in _CRITICAL_GA4_TAGS if k not in found_set]

    print(f"\nPPT filled successfully: {output_path}")
    print(f"Placeholders replaced: {len(placeholders_found)}")
    if missing_critical:
        print(
            "WARNING — these GA4 total tags were NOT found in the template "
            "(table may show stale/wrong numbers for those columns):"
        )
        for k in missing_critical:
            print(f"  - {{{{{k}}}}}")
    if placeholders_missing:
        print(f"WARNING — {len(placeholders_missing)} placeholder fragments remain:")
        for m in list(dict.fromkeys(placeholders_missing))[:15]:
            print(f"  - {m}")

    return {
        "found": placeholders_found,
        "missing": placeholders_missing,
        "unknown_template_tags": unknown_tags,
        "missing_critical_tags": missing_critical,
    }


def run_ppt_job(
    cfg,
    start_date,
    end_date,
    month_label,
    progress_cb=None,
    prev_start_date=None,
    prev_end_date=None,
):
    """progress_cb receives dicts: kind \"log\" | \"download\", optional percent, message."""

    def emit(payload):
        if progress_cb:
            progress_cb(payload)

    drive = get_drive_service()

    with tempfile.TemporaryDirectory(ignore_cleanup_errors=True) as tmp_dir:
        template_path = os.path.join(tmp_dir, "template.pptx")
        output_path = os.path.join(tmp_dir, "output.pptx")

        print("Downloading template from Google Drive...")
        emit({"kind": "log", "message": "Downloading template from Google Drive..."})
        download_file(drive, cfg["template_drive_id"], template_path, progress_cb=emit)

        print("Fetching data from Supabase/GA4...")
        emit({"kind": "log", "message": "Fetching data from Supabase/GA4..."})
        prev_start = (prev_start_date or "").strip() or None
        prev_end = (prev_end_date or "").strip() or None
        data = build_full_data(
            cfg["customer_id"],
            start_date,
            end_date,
            ga4_property_id=cfg.get("ga4_property_id"),
            customer_name_fallback=cfg.get("client_name"),
            prev_start_date=prev_start,
            prev_end_date=prev_end,
            client_id=cfg.get("client_id"),
        )
        emit({
            "kind": "log",
            "message": (
                f"GA4 totals: views={data.get('ga4_total_views')} "
                f"sessions={data.get('ga4_total_sessions')} "
                f"users={data.get('ga4_total_users')} "
                f"bounce={data.get('ga4_total_bounce')}"
            ),
        })
        print(f"Data fetched: {len(data)} values")
        emit({"kind": "log", "message": f"Data fetched: {len(data)} values"})

        print("Filling the template...")
        emit({"kind": "log", "message": "Filling the template..."})
        result = fill_presentation(data, template_path, output_path)
        print(f"\nPPT filled successfully: {output_path}")
        print(f"Placeholders replaced: {len(result['found'])}")

        missing_critical = result.get("missing_critical_tags") or []
        if missing_critical:
            msg = (
                "WARNING: Template missing tags "
                + ", ".join(f"{{{{{k}}}}}" for k in missing_critical[:5])
                + " — Users/Sessions columns may show old hard-coded numbers. "
                "Fix the .pptx template on Google Drive."
            )
            emit({"kind": "log", "message": msg})

        unknown = result.get("unknown_template_tags") or []
        if unknown:
            emit({
                "kind": "log",
                "message": (
                    "WARNING: Unknown template tags: "
                    + ", ".join(f"{{{{{t}}}}}" for t in unknown[:8])
                    + (" ..." if len(unknown) > 8 else "")
                ),
            })

        emit({"kind": "log", "message": f"PPT filled successfully: {output_path}"})
        emit({"kind": "log", "message": f"Placeholders replaced: {len(result['found'])}"})

        output_filename = f"{cfg['client_name']}_{month_label}_filled.pptx"
        print("Uploading and overwriting file on Google Drive...")

        uploaded = update_file(
            drive,
            cfg["output_file_drive_id"],
            output_path,
            output_filename,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            progress_cb=emit,
        )

        drive_url = uploaded.get("webViewLink")
        emit({"kind": "log", "message": "\n=== RESULT ==="})
        emit({"kind": "log", "message": f"File: {output_filename}"})
        emit({"kind": "log", "message": f"URL:  {drive_url}"})

        return {
            "filename": output_filename,
            "drive_url": drive_url,
            "placeholders_found": len(result["found"]),
            "missing_critical_tags": missing_critical,
            "unknown_template_tags": unknown,
        }


if __name__ == "__main__":
    import json

    with open("config/zoomers_rv.json") as f:
        cfg = json.load(f)

    result = run_ppt_job(cfg, "2026-04-01", "2026-04-30", "April_2026")

    print("\n=== RESULT ===")
    print(f"File: {result['filename']}")
    print(f"URL:  {result['drive_url']}")
